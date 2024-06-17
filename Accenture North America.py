#!/usr/bin/env python
# coding: utf-8

# In[5]:


import pandas as pd

# Load the datasets
content_df = pd.read_csv('Content.csv')
reactions_df = pd.read_csv('Reactions.csv')
reaction_types_df = pd.read_csv('ReactionTypes.csv')

# Display the first few rows of each dataframe to understand their structure
content_df.head(), reactions_df.head(), reaction_types_df.head()

# Clean the Content dataset
# Remove rows with missing values
content_df_cleaned = content_df.dropna()

# Remove columns that are not relevant
content_df_cleaned = content_df_cleaned.drop(columns=['Unnamed: 0'])

# Clean the Reactions dataset
# Remove rows with missing values
reactions_df_cleaned = reactions_df.dropna()

# Remove columns that are not relevant
reactions_df_cleaned = reactions_df_cleaned.drop(columns=['Unnamed: 0'])

# Convert the Datetime column to datetime format
reactions_df_cleaned['Datetime'] = pd.to_datetime(reactions_df_cleaned['Datetime'])

# Clean the ReactionTypes dataset
# Remove rows with missing values
reaction_types_df_cleaned = reaction_types_df.dropna()

# Remove columns that are not relevant
reaction_types_df_cleaned = reaction_types_df_cleaned.drop(columns=['Unnamed: 0'])

# Function to display the cleaned dataframes to the user
def display_dataframe_to_user(name, dataframe):
    dataframe.to_csv(f"{name}.csv", index=False)

# Display cleaned dataframes to the user
display_dataframe_to_user("Cleaned_Content_Data", content_df_cleaned)
display_dataframe_to_user("Cleaned_Reactions_Data", reactions_df_cleaned)
display_dataframe_to_user("Cleaned_Reaction_Types_Data", reaction_types_df_cleaned)

content_df_cleaned.head(), reactions_df_cleaned.head(), reaction_types_df_cleaned.head()


# In[9]:


# Merge the dataframes together
merged_df = reactions_df.merge(content_df[['Content ID', 'Category']], on='Content ID', how='left')
merged_df = merged_df.merge(reaction_types_df[['Type', 'Score']], on='Type', how='left')

# Display the first few rows of the merged dataframe
merged_df.head()


# In[10]:


# Calculate the total scores for each category
category_scores = merged_df.groupby('Category')['Score'].sum().reset_index()

# Sort the categories by total score in descending order and get the top 5
top_5_categories = category_scores.sort_values(by='Score', ascending=False).head(5)

# Display the top 5 categories
top_5_categories


# In[14]:


# Save the cleaned dataset and the top 5 categories to a new file
cleaned_and_top5_df = merged_df.copy()

# Add a column indicating if the category is in the top 5
cleaned_and_top5_df['Top 5 Category'] = cleaned_and_top5_df['Category'].apply(lambda x: x if x in top_5_categories['Category'].values else 'Other')

# Save to CSV
output_file = 'Cleaned_and_Top5_Categories.csv'
cleaned_and_top5_df.to_csv(output_file, index=False)


output_file


# In[18]:


import matplotlib.pyplot as plt
from pptx.util import Inches
from pptx import Presentation
import pandas as pd

# Load the datasets
reactions_df = pd.read_csv('Cleaned_Reactions_Data.csv')
reaction_types_df = pd.read_csv('Cleaned_Reaction_Types_Data.csv')
content_df = pd.read_csv('Cleaned_Content_Data.csv')

# Merge the dataframes together
merged_df = reactions_df.merge(content_df[['Content ID', 'Category']], on='Content ID', how='left')
merged_df = merged_df.merge(reaction_types_df[['Type', 'Score']], on='Type', how='left')

# Calculate the total scores for each category
category_scores = merged_df.groupby('Category')['Score'].sum().reset_index()

# Sort the categories by total score in descending order and get the top 5
top_5_categories = category_scores.sort_values(by='Score', ascending=False).head(5)

# Create bar plot for the top 5 categories
plt.figure(figsize=(10, 6))
plt.bar(top_5_categories['Category'], top_5_categories['Score'], color='skyblue')
plt.xlabel('Category')
plt.ylabel('Total Score')
plt.title('Top 5 Performing Categories')
plt.tight_layout()
plt.savefig('top_5_categories_chart.png')
plt.close()

# Load the provided PowerPoint template
ppt_template_path = 'Data Analytics template - Task 3_final.pptx'
presentation = Presentation(ppt_template_path)

# Define content for each slide
slides_content = [
    {"title": "Top Performing Categories Analysis", "content": "Insights from Data Analytics\n[Your Name]\n[Today's Date]"},
    {"title": "Agenda", "content": (
        "Today's Agenda\n"
        "1. Project Recap\n"
        "2. Problem\n"
        "3. The Analytics Team\n"
        "4. Process\n"
        "5. Insights\n"
        "6. Summary\n"
        "7. Questions"
    )},
    {"title": "Project Recap", "content": (
        "Objective: Identify the top 5 performing categories\n"
        "Data Sources: Reactions, Content, Reaction Types"
    )},
    {"title": "Problem", "content": (
        "How to determine the top 5 performing categories based on reaction scores?"
    )},
    {"title": "The Analytics Team", "content": (
        "Andrew Fleming - Chief Technical Architect\n"
        "Marcus Rompton - Senior Principal\n"
        "[Your Name] - [Your Title]"
    )},
    {"title": "Process", "content": (
        "1. Data Collection\n"
        "2. Data Cleaning\n"
        "3. Data Merging\n"
        "4. Data Analysis\n"
        "5. Insights Generation"
    )},
    {"title": "Insights", "content": (
        "Top 5 Performing Categories:\n"
        "1. Travel\n"
        "2. Science\n"
        "3. Healthy Eating\n"
        "4. Animals\n"
        "5. Cooking"
    )},
    {"title": "Summary", "content": (
        "Key Takeaways:\n"
        "- Travel and Science are the top categories\n"
        "- Detailed performance metrics"
    )},
    {"title": "Questions", "content": "Any Questions?\nThank you!"},
]

# Add content to the slides
for slide_content in slides_content:
    slide_layout = presentation.slide_layouts[1]  # Assuming layout 1 is a title and content layout
    slide = presentation.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = slide_content["title"]
    content_shape = slide.placeholders[1]
    content_shape.text = slide_content["content"]

# Add the visualization to the Insights slide
insights_slide = presentation.slides[7]  
img_path = 'top_5_categories_chart.png'
insights_slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8), height=Inches(4.5))

# Save the modified presentation
output_ppt_path = 'Top_Performing_Categories_Analysis_Presentation_with_Charts.pptx'
presentation.save(output_ppt_path)

output_ppt_path

